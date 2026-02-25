import streamlit as st
import os
import json
from groq import Groq
from fpdf import FPDF
from pptx import Presentation
from dotenv import load_dotenv

# --- 1. SETUP GROQ ---
load_dotenv()
API_KEY = os.getenv("GROQ_API_KEY")

if API_KEY:
    client = Groq(api_key=API_KEY)
else:
    st.error("Groq API Key missing! Update your .env file with GROQ_API_KEY.")
    st.stop()

st.set_page_config(page_title="Apex Suite", layout="wide")

# --- SESSION STATE ---
if "app_mode" not in st.session_state: st.session_state.app_mode = "Hub"
if "step" not in st.session_state: st.session_state.step = 1

# Resume Fields
resume_fields = ["name", "email", "phone", "linkedin", "github", "edu_rough", "hobbies_rough", "exp_rough", "proj_rough", "ai_edu", "ai_hobbies"]
for f in resume_fields:
    if f not in st.session_state: st.session_state[f] = ""
if "ai_exp_options" not in st.session_state: st.session_state.ai_exp_options = []
if "ai_proj_options" not in st.session_state: st.session_state.ai_proj_options = []

# Portfolio Fields
if "port_name" not in st.session_state: st.session_state.port_name = ""
if "port_tagline" not in st.session_state: st.session_state.port_tagline = ""
if "port_about" not in st.session_state: st.session_state.port_about = ""
if "port_projects" not in st.session_state: st.session_state.port_projects = ""
if "port_slides" not in st.session_state: st.session_state.port_slides = []

def reset_wizard():
    st.session_state.step = 1

# ==========================================
# UI: THE HUB
# ==========================================
if st.session_state.app_mode == "Hub":
    st.markdown("<h1 style='text-align: center; font-size: 3em;'>Welcome to the Apex Suite</h1>", unsafe_allow_html=True)
    st.markdown("<p style='text-align: center; color: gray;'>Select your project below to get started.</p><br><br>", unsafe_allow_html=True)
    
    col1, col2, col3, col4 = st.columns([1, 4, 4, 1])
    
    with col2:
        st.markdown("<h1 style='text-align: center; font-size: 5em;'>📄</h1>", unsafe_allow_html=True)
        st.markdown("<h3 style='text-align: center;'>Resume Builder</h3>", unsafe_allow_html=True)
        if st.button("Launch Resume Track", use_container_width=True, type="primary"):
            st.session_state.app_mode = "Resume"
            reset_wizard()
            st.rerun()
            
    with col3:
        st.markdown("<h1 style='text-align: center; font-size: 5em;'>🖼️</h1>", unsafe_allow_html=True)
        st.markdown("<h3 style='text-align: center;'>Portfolio Builder (PPTX)</h3>", unsafe_allow_html=True)
        if st.button("Launch Portfolio Track", use_container_width=True, type="primary"):
            st.session_state.app_mode = "Portfolio"
            reset_wizard()
            st.rerun()

# ==========================================
# UI: PATH A - RESUME BUILDER (UNTOUCHED)
# ==========================================
elif st.session_state.app_mode == "Resume":
    if st.sidebar.button("← Back to Hub"):
        st.session_state.app_mode = "Hub"
        st.rerun()
        
    st.title("Apex Resume Builder")
    st.progress(st.session_state.step / 6)

    def process_all_ai_content():
        mega_prompt = f"""
        You are an expert executive resume writer. Process the user's rough notes into professional resume content.
        User Notes:
        Education: {st.session_state.edu_rough}
        Skills: {st.session_state.hobbies_rough}
        Experience: {st.session_state.exp_rough}
        Projects: {st.session_state.proj_rough}
        Formatting Rules:
        - Wrap metrics and key terms in <b> tags.
        - Wrap tools and software in <i> tags.
        - "edu" and "skills" must be single formatted strings.
        - "exp" and "proj" must be lists containing exactly 4 detailed bullet points each.
        You must output valid JSON.
        """
        try:
            chat_completion = client.chat.completions.create(
                messages=[
                    {"role": "system", "content": "You are a helpful resume assistant that outputs JSON.\nThe JSON object must use the schema: {'edu': 'str', 'skills': 'str', 'exp': ['str', 'str', 'str', 'str'], 'proj': ['str', 'str', 'str', 'str']}"},
                    {"role": "user", "content": mega_prompt}
                ],
                model="llama-3.3-70b-versatile",
                response_format={"type": "json_object"},
            )
            data = json.loads(chat_completion.choices[0].message.content)
            st.session_state.ai_edu = data.get("edu", "Education details not provided.")
            st.session_state.ai_hobbies = data.get("skills", "Skill details not provided.")
            st.session_state.ai_exp_options = data.get("exp", [])
            st.session_state.ai_proj_options = data.get("proj", [])
            return True
        except Exception as e:
            st.error(f"Groq API Error: {e}")
            return False

    def create_pdf(sel_exp, sel_proj):
        pdf = FPDF()
        pdf.add_page()
        def safe(t): return t.encode('latin-1', 'replace').decode('latin-1')
        pdf.set_font("Helvetica", "B", 24)
        pdf.cell(0, 10, safe(st.session_state.name.upper()), align="C", ln=True)
        pdf.set_font("Helvetica", "", 11)
        pdf.cell(0, 6, safe(f"{st.session_state.email} | {st.session_state.phone}"), align="C", ln=True)
        pdf.ln(8)

        def draw_sec(title, items):
            pdf.set_font("Helvetica", "B", 14)
            pdf.cell(0, 8, safe(title), ln=True)
            pdf.line(10, pdf.get_y(), 200, pdf.get_y()); pdf.ln(3)
            pdf.set_font("Helvetica", size=11)
            if isinstance(items, list):
                for i in items:
                    clean_item = i.replace('<b>','').replace('</b>','').replace('<i>','').replace('</i>','')
                    pdf.multi_cell(0, 6, safe(f"- {clean_item}")); pdf.ln(1)
            else:
                clean_item = items.replace('<b>','').replace('</b>','').replace('<i>','').replace('</i>','')
                pdf.write_html(safe(clean_item)); pdf.ln(5)

        if st.session_state.ai_edu: draw_sec("EDUCATION", st.session_state.ai_edu)
        if sel_exp: draw_sec("PROFESSIONAL EXPERIENCE", sel_exp)
        if sel_proj: draw_sec("TECHNICAL PROJECTS", sel_proj)
        if st.session_state.ai_hobbies: draw_sec("SKILLS & HOBBIES", st.session_state.ai_hobbies)
        path = "Apex_Resume.pdf"
        pdf.output(path)
        return path

    if st.session_state.step == 1:
        st.header("Step 1: Contact")
        st.session_state.name = st.text_input("Name", st.session_state.name)
        st.session_state.email = st.text_input("Email", st.session_state.email)
        st.session_state.phone = st.text_input("Phone", st.session_state.phone)
        if st.button("Next"): st.session_state.step = 2; st.rerun()

    elif st.session_state.step == 2:
        st.header("Step 2: Links")
        st.session_state.linkedin = st.text_input("LinkedIn", st.session_state.linkedin)
        st.session_state.github = st.text_input("GitHub", st.session_state.github)
        c1, c2 = st.columns(2)
        if c1.button("Back"): st.session_state.step = 1; st.rerun()
        if c2.button("Next"): st.session_state.step = 3; st.rerun()

    elif st.session_state.step == 3:
        st.header("Step 3: Edu & Skills")
        st.session_state.edu_rough = st.text_area("Education", st.session_state.edu_rough)
        st.session_state.hobbies_rough = st.text_area("Skills", st.session_state.hobbies_rough)
        c1, c2 = st.columns(2)
        if c1.button("Back"): st.session_state.step = 2; st.rerun()
        if c2.button("Next"): st.session_state.step = 4; st.rerun()

    elif st.session_state.step == 4:
        st.header("Step 4: Experience")
        st.session_state.exp_rough = st.text_area("Experience", st.session_state.exp_rough, height=150)
        c1, c2 = st.columns(2)
        if c1.button("Back"): st.session_state.step = 3; st.rerun()
        if c2.button("Next"): st.session_state.step = 5; st.rerun()

    elif st.session_state.step == 5:
        st.header("Step 5: Projects")
        st.session_state.proj_rough = st.text_area("Projects", st.session_state.proj_rough, height=150)
        c1, c2 = st.columns(2)
        if c1.button("Back"): st.session_state.step = 4; st.rerun()
        if c2.button("Generate AI Content", type="primary"):
            with st.spinner("Groq Llama-3.3 is writing your resume..."):
                if process_all_ai_content():
                    st.session_state.step = 6
                    st.rerun()

    elif st.session_state.step == 6:
        st.header("Step 6: Review & Finalize")
        st.markdown(f"**Edu:** {st.session_state.ai_edu}", unsafe_allow_html=True)
        st.markdown(f"**Skills:** {st.session_state.ai_hobbies}", unsafe_allow_html=True)
        st.divider()
        e = [p for i, p in enumerate(st.session_state.ai_exp_options) if st.checkbox(p, value=True, key=f"e{i}")]
        st.divider()
        p = [p for i, p in enumerate(st.session_state.ai_proj_options) if st.checkbox(p, value=True, key=f"p{i}")]
        if st.button("Download PDF"):
            path = create_pdf(e, p)
            with open(path, "rb") as f: st.download_button("Download", f, file_name="Apex_Resume.pdf")

# ==========================================
# UI: PATH B - PORTFOLIO BUILDER
# ==========================================
elif st.session_state.app_mode == "Portfolio":
    if st.sidebar.button("← Back to Hub"):
        st.session_state.app_mode = "Hub"
        st.rerun()

    st.title("Apex Portfolio Builder")
    st.progress(st.session_state.step / 4)
    
    def generate_portfolio_slides():
        prompt = f"""
        Act as a professional presentation copywriter. I need content for a 6-slide portfolio deck.
        Name: {st.session_state.port_name}
        Tagline: {st.session_state.port_tagline}
        About Me Notes: {st.session_state.port_about}
        Project Notes: {st.session_state.port_projects}
        
        Create exactly 6 slides. Return valid JSON only in this format:
        {{"slides": [ {{"title": "Slide Title", "content": "Slide bullet points or paragraph..."}} ]}}
        """
        try:
            chat = client.chat.completions.create(
                messages=[
                    {"role": "system", "content": "You are a JSON-only API. Schema: {'slides': [{'title': 'str', 'content': 'str'}]}"},
                    {"role": "user", "content": prompt}
                ],
                model="llama-3.3-70b-versatile",
                response_format={"type": "json_object"},
            )
            data = json.loads(chat.choices[0].message.content)
            st.session_state.port_slides = data.get("slides", [])
            return True
        except Exception as e:
            st.error(f"Groq API Error: {e}")
            return False

    def create_pptx(slides_data):
        prs = Presentation()
        
        # 1. Title Slide (Layout 0)
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = st.session_state.port_name
        title_slide.placeholders[1].text = st.session_state.port_tagline
        
        # 2. Add AI Content Slides (Layout 1)
        for slide in slides_data:
            s = prs.slides.add_slide(prs.slide_layouts[1])
            s.shapes.title.text = slide["title"]
            s.placeholders[1].text = slide["content"]
            
        path = "Apex_Portfolio.pptx"
        prs.save(path)
        return path

    if st.session_state.step == 1:
        st.header("Step 1: The Basics")
        st.session_state.port_name = st.text_input("Your Name", st.session_state.port_name)
        st.session_state.port_tagline = st.text_input("Your Professional Tagline (e.g., Data Scientist | AI Architect)", st.session_state.port_tagline)
        if st.button("Next"): st.session_state.step = 2; st.rerun()

    elif st.session_state.step == 2:
        st.header("Step 2: About Me & Projects")
        st.session_state.port_about = st.text_area("About Me Notes (What drives you?)", st.session_state.port_about, height=100)
        st.session_state.port_projects = st.text_area("Project Case Studies (Dump your rough notes here)", st.session_state.port_projects, height=200)
        c1, c2 = st.columns(2)
        if c1.button("Back"): st.session_state.step = 1; st.rerun()
        if c2.button("Generate Slides", type="primary"):
            with st.spinner("Groq is architecting your PowerPoint deck..."):
                if generate_portfolio_slides():
                    st.session_state.step = 3
                    st.rerun()

    elif st.session_state.step == 3:
        st.header("Step 3: Review Slides & Download")
        st.success("PowerPoint content generated! Review the outline below.")
        
        for idx, slide in enumerate(st.session_state.port_slides):
            with st.expander(f"Slide {idx + 2}: {slide['title']}", expanded=True):
                st.write(slide['content'])
                
        st.divider()
        if st.button("Download PowerPoint (.pptx)", type="primary"):
            path = create_pptx(st.session_state.port_slides)
            with open(path, "rb") as f:
                st.download_button(label="Click to save your PPTX", data=f, file_name="Apex_Portfolio.pptx")
